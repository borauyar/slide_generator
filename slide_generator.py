import os
import sys
import re
import pathlib
from googlesearch import search
from google_images_search import GoogleImagesSearch
from pptx import Presentation
from pptx.util import Inches

def parse_markdown(file_path):
    with open(file_path, 'r') as f:
        content = f.read()

    slides = content.split('#')[1:]
    slide_data = []

    for slide in slides:
        lines = slide.strip().split('\n')
        title = lines[0].strip()
        bullet_points = [line.strip().replace('-', '').strip() for line in lines[1:]]
        slide_data.append({'title': title, 'bullet_points': bullet_points})

    return slide_data

def download_image(keyword, api_key, cx):
    gis = GoogleImagesSearch(api_key, cx)

    _search_params = {
        'q': keyword,
        'num': 1,
        'imgSize': 'large',
        'safe': 'off',
        'fileType': 'jpg|gif|png',
        'rights': 'cc_publicdomain|cc_attribute|cc_sharealike|cc_noncommercial|cc_nonderived',
    }

    gis.search(search_params=_search_params)
    if not gis.results():
        return None

    image = gis.results()[0]
    img_dir = f'./downloads/{keyword.replace(" ", "_")}/'
    img_name = f'image.jpg'
    pathlib.Path(img_dir).mkdir(parents=True, exist_ok=True)

    # Save the image to the desired location
    with open(os.path.join(img_dir, img_name), 'wb') as f:
        f.write(image.get_raw_data())

    return os.path.join(img_dir, img_name)

def create_presentation(slide_data, api_key, cx):
    prs = Presentation()

    for slide_info in slide_data:
        title = slide_info['title']
        bullet_points = slide_info['bullet_points']

        image_path = None
        for keyword in bullet_points:
            image_path = download_image(keyword, api_key, cx)
            if image_path:
                break

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title

        if image_path:
            img_left = Inches(4.5)
            img_top = Inches(1.5)
            img_width = Inches(4)
            img_height = Inches(3)
            slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)

        bullet_shape = slide.shapes.placeholders[1].text_frame
        for point in bullet_points:
            p = bullet_shape.add_paragraph()
            p.text = point
            p.level = 0

    prs.save('presentation.pptx')

def main():
    if len(sys.argv) != 4:
        print("Usage: python script.py markdown_file_path api_key cx")
        sys.exit(1)

    markdown_file_path = sys.argv[1]
    api_key = sys.argv[2]
    cx = sys.argv[3]

    slide_data = parse_markdown(markdown_file_path)
    create_presentation(slide_data, api_key, cx)

if __name__ == '__main__':
    main()


